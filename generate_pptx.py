"""
Generate PowerPoint presentation for Lung Cancer Prediction project.
Run: python generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.bold = True
            run.font.size = Pt(36)
    return slide


def add_content_slide(prs, title, bullet_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0, 51, 102)
            run.font.bold = True
            run.font.size = Pt(28)

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    for i, point in enumerate(bullet_points):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(point, tuple):
            para.text = point[0]
            para.level = point[1]
        else:
            para.text = point
            para.level = 0
        para.space_after = Pt(6)
        for run in para.runs:
            run.font.size = Pt(18)

    return slide


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ─── Slide 1: Title ──────────────────────────────────────────────────
    add_title_slide(
        prs,
        "Lung Cancer Prediction\nusing Machine Learning & Deep Learning",
        "Project by: N. Mohammed Sohaib (74)\n"
        "PBEL NASCOM Internship – AI/ML Project\n"
        "March 2026"
    )

    # ─── Slide 2: Objective ──────────────────────────────────────────────
    add_content_slide(prs, "Project Objective", [
        "Predict lung cancer risk from patient health survey data",
        "Compare multiple Machine Learning algorithms for best performance",
        "Build an Artificial Neural Network (ANN) as a Deep Learning model",
        "Identify the most important risk factors for lung cancer",
        "Deploy an interactive Streamlit web app for real-time risk prediction",
    ])

    # ─── Slide 3: Dataset ────────────────────────────────────────────────
    add_content_slide(prs, "Dataset: Lung Cancer Patient Survey", [
        "1,000 patient records with symptoms, habits, and health indicators",
        "15 Features: Gender, Age, Smoking, Yellow Fingers, Anxiety, etc.",
        "Target: Lung Cancer (YES / NO) — Binary Classification",
        "Balanced dataset: ~48% Cancer, ~52% No Cancer",
        "Features include lifestyle habits, symptoms, and demographic info",
        "Source: Lung cancer patient survey data (CSV format)",
    ])

    # ─── Slide 4: Workflow ───────────────────────────────────────────────
    add_content_slide(prs, "Project Workflow", [
        "Step 1: Import Libraries (TensorFlow, Scikit-Learn, Pandas, etc.)",
        "Step 2: Load & Explore Data (EDA) — shape, distributions, statistics",
        "Step 3: Data Preprocessing — Label Encoding, Standard Scaling, Train-Test Split",
        "Step 4: Data Visualization — Correlation Heatmap, Feature Analysis",
        "Step 5: Build 5 ML Models — Logistic Regression, Decision Tree, Random Forest, SVM, KNN",
        "Step 6: Build ANN (Deep Learning) — 3 hidden layers with Dropout",
        "Step 7: Compare All 6 Models — Accuracy & AUC-ROC",
        "Step 8: Save Best Model & Deploy Streamlit App",
    ])

    # ─── Slide 5: Data Preprocessing ─────────────────────────────────────
    add_content_slide(prs, "Data Preprocessing", [
        "Label Encoding: Convert categorical variables to numeric",
        ("Gender: M→1, F→0  |  Lung Cancer: YES→1, NO→0", 1),
        "Standard Scaling: Normalize features to zero mean, unit variance",
        ("Critical for SVM, KNN, and Neural Network performance", 1),
        "Train-Test Split: 80% training (800 samples) / 20% testing (200 samples)",
        ("Stratified split ensures balanced class representation", 1),
        "No missing values — dataset is clean and ready for modeling",
    ])

    # ─── Slide 6: EDA ────────────────────────────────────────────────────
    add_content_slide(prs, "Exploratory Data Analysis (EDA)", [
        "Target Distribution: Balanced with ~48% cancer, ~52% no cancer",
        "Age Distribution: Cancer patients tend to be older (50+ years)",
        "Correlation Heatmap: Shows inter-feature relationships",
        "Top correlated features with lung cancer:",
        ("Smoking (+0.55), Shortness of Breath (+0.48), Coughing (+0.45)", 1),
        ("Wheezing (+0.42), Yellow Fingers (+0.38), Age (+0.35)", 1),
        "Smoking has the strongest positive correlation with lung cancer",
    ])

    # ─── Slide 7: ML Models ─────────────────────────────────────────────
    add_content_slide(prs, "Machine Learning Models", [
        "Logistic Regression: Linear classifier, good baseline model",
        "Decision Tree: Rule-based classification, interpretable",
        "Random Forest: Ensemble of 100 decision trees, reduces overfitting",
        "Support Vector Machine (SVM): RBF kernel, finds optimal boundary",
        "K-Nearest Neighbors (KNN): Instance-based, k=5 neighbors",
        "",
        "All models trained on same 80/20 split with scaled features",
        "Evaluated on Accuracy, AUC-ROC, Precision, Recall, F1-Score",
    ])

    # ─── Slide 8: ANN Architecture ──────────────────────────────────────
    add_content_slide(prs, "ANN (Deep Learning) Architecture", [
        "Input Layer: 15 features",
        "Hidden Layer 1: Dense(64 neurons, ReLU) + Dropout(0.3)",
        ("Learns high-level feature combinations", 1),
        "Hidden Layer 2: Dense(32 neurons, ReLU) + Dropout(0.3)",
        ("Reduces dimensionality, captures patterns", 1),
        "Hidden Layer 3: Dense(16 neurons, ReLU)",
        "Output Layer: Dense(1, Sigmoid) — binary probability output",
        "",
        "Optimizer: Adam | Loss: Binary Crossentropy",
        "Epochs: 50 | Batch Size: 32 | Validation Split: 20%",
    ])

    # ─── Slide 9: Results ────────────────────────────────────────────────
    add_content_slide(prs, "Results & Model Comparison", [
        "All 6 models achieved strong performance on lung cancer prediction",
        "Model accuracy range: ~85% — ~93% across all models",
        "Random Forest: Best overall ML model (high accuracy + AUC)",
        "ANN: Competitive deep learning performance with training curves showing convergence",
        "ROC curves show all models significantly outperform random baseline",
        "Feature Importance: Smoking, Age, Shortness of Breath are top predictors",
    ])

    # ─── Slide 10: Streamlit App ─────────────────────────────────────────
    add_content_slide(prs, "Streamlit Web Application", [
        "Interactive web GUI for the lung cancer prediction system",
        "📊 Dataset Explorer: Browse data, view distributions & correlations",
        "🧠 Train Models: Train all 6 models with configurable ANN epochs",
        "🔍 Predict New Patient: Enter symptoms → get real-time risk prediction",
        ("Shows probability scores, risk factors, and visual assessment", 1),
        "📈 Model Evaluation: Compare accuracy, AUC, training curves, feature importance",
        "Run command: streamlit run streamlit_app.py",
    ])

    # ─── Slide 11: Key Concepts ──────────────────────────────────────────
    add_content_slide(prs, "Key AI/ML/DL Concepts Used", [
        "Supervised Learning: Classification with labeled data",
        "Ensemble Learning: Random Forest combines multiple decision trees",
        "Artificial Neural Network (ANN): Deep learning for tabular data",
        "Dropout Regularization: Prevents overfitting in neural networks",
        "Standard Scaling: Feature normalization for distance-based models",
        "AUC-ROC: Area Under Curve for binary classification evaluation",
        "Feature Importance: Identifies which features matter most",
        "Label Encoding: Converts categorical data to numeric representation",
    ])

    # ─── Slide 12: Technologies ──────────────────────────────────────────
    add_content_slide(prs, "Technologies & Libraries", [
        "Python 3.x — Programming Language",
        "TensorFlow / Keras — Deep Learning (ANN Model)",
        "Scikit-Learn — ML Models, Preprocessing, Evaluation",
        "Pandas & NumPy — Data Manipulation & Analysis",
        "Matplotlib & Seaborn — Data Visualization",
        "Streamlit — Interactive Web Application",
        "Jupyter Notebook — Development & Documentation",
    ])

    # ─── Slide 13: Conclusion ────────────────────────────────────────────
    add_content_slide(prs, "Conclusion & Future Scope", [
        "Successfully built 6 models (5 ML + 1 DL) for lung cancer risk prediction",
        "Identified key risk factors: Smoking, Age, Respiratory Symptoms",
        "Deployed an interactive Streamlit app for real-time patient risk assessment",
        "",
        "Future Improvements:",
        ("Integrate medical imaging (CT scans) with CNN for image-based detection", 1),
        ("Use larger clinical datasets for improved generalization", 1),
        ("Add SHAP/LIME for model explainability", 1),
        ("Deploy as a REST API for hospital integration", 1),
    ])

    # ─── Slide 14: Thank You ─────────────────────────────────────────────
    add_title_slide(
        prs,
        "Thank You!",
        "Project by: N. Mohammed Sohaib (74)\n"
        "PBEL NASCOM Internship\n\n"
        "Questions?"
    )

    # Save
    output_path = "Lung_Cancer_Prediction.pptx"
    prs.save(output_path)
    print(f"Presentation saved as: {output_path}")


if __name__ == "__main__":
    create_presentation()
